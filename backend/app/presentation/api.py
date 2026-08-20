from __future__ import annotations

import io
import re
from datetime import date, datetime
from typing import Annotated
from uuid import UUID

from fastapi import APIRouter, Depends, Form, Header, HTTPException, Response, status
from pydantic import BaseModel, Field

from app.application.auth_service import AuthService
from app.application.contracts import IdempotencyConflictError, SubmitCommand, UpdateCommand
from app.application.services import (
    ApproveProjectUpdate,
    CheckExistingReport,
    GetProjectUpdate,
    ListProjectUpdates,
    SubmitProjectUpdate,
    UpdateProjectUpdate,
)
from app.domain import ProjectUpdate
from app.domain.exceptions import DomainValidationError

_TAG_RE = re.compile(r"<[^>]+>")


def _strip_html(html: str) -> str:
    return _TAG_RE.sub("", html or "").replace("&nbsp;", " ").strip()


class UpdateProjectUpdateRequest(BaseModel):
    start_of_week: date
    end_of_week: date
    team_project: str = Field(min_length=1, max_length=300)
    achievements: str = Field(min_length=1, max_length=5000)
    initiatives: str = Field(min_length=1, max_length=5000)
    next_weeks_plan: str = Field(min_length=1, max_length=5000)


def serialize(update: ProjectUpdate, owner_username: str | None = None) -> dict[str, object]:
    return {
        "id": str(update.id),
        "user_id": str(update.user_id) if update.user_id else None,
        "owner_username": owner_username,
        "start_of_week": update.start_of_week,
        "end_of_week": update.end_of_week,
        "team_project": update.team_project,
        "achievements": update.achievements,
        "initiatives": update.initiatives,
        "next_weeks_plan": update.next_weeks_plan,
        "created_at": update.created_at,
        "updated_at": update.updated_at,
        "approval_status": update.approval_status,
        "approved_by_id": str(update.approved_by_id) if update.approved_by_id else None,
        "approved_at": update.approved_at,
    }


def create_project_update_router(
    submit: SubmitProjectUpdate,
    get_update: GetProjectUpdate,
    list_updates_service: ListProjectUpdates,
    update_service: UpdateProjectUpdate,
    check_existing: CheckExistingReport,
    approve_service: ApproveProjectUpdate,
    get_current_user,
    auth_service: AuthService,
) -> APIRouter:
    router = APIRouter(prefix="/api/v1/project-updates", tags=["project updates"])

    @router.post("", status_code=status.HTTP_201_CREATED)
    def create_update(
        response: Response,
        idempotency_key: Annotated[str, Header(alias="Idempotency-Key", min_length=1, max_length=128)],
        start_of_week: Annotated[date, Form()],
        end_of_week: Annotated[date, Form()],
        team_project: Annotated[str, Form(min_length=1, max_length=300)],
        achievements: Annotated[str, Form(min_length=1, max_length=5000)],
        initiatives: Annotated[str, Form(min_length=1, max_length=5000)],
        next_weeks_plan: Annotated[str, Form(min_length=1, max_length=5000)],
        current_user: dict = Depends(get_current_user),
    ) -> dict[str, object]:
        user_role = current_user.get("role")
        if user_role not in ("TEAM_LEAD", "APP_ADMIN"):
            raise HTTPException(status.HTTP_403_FORBIDDEN, "Only Team Leads can create reports.")
        current_user_id = UUID(current_user["sub"])
        try:
            result = submit.execute(
                SubmitCommand(
                    idempotency_key=idempotency_key,
                    start_of_week=start_of_week,
                    end_of_week=end_of_week,
                    team_project=team_project,
                    achievements=achievements,
                    initiatives=initiatives,
                    next_weeks_plan=next_weeks_plan,
                ),
                user_id=current_user_id,
            )
        except IdempotencyConflictError as error:
            raise HTTPException(status.HTTP_409_CONFLICT, str(error)) from error
        except DomainValidationError as error:
            raise HTTPException(status.HTTP_422_UNPROCESSABLE_ENTITY, str(error)) from error
        if result.replayed:
            response.status_code = status.HTTP_200_OK
            response.headers["Idempotent-Replayed"] = "true"
        return serialize(result.update)

    # --- Export endpoints (DU_HEAD only) must be registered before /{update_id} ---

    @router.get("/export/excel")
    def export_excel(current_user: dict = Depends(get_current_user)) -> Response:
        import openpyxl  # type: ignore[import]
        import openpyxl.styles  # type: ignore[import]

        if current_user.get("role") != "DU_HEAD":
            raise HTTPException(status.HTTP_403_FORBIDDEN, "Only DU Head can export reports.")
        updates = list_updates_service.execute()
        owner_ids = {u.user_id for u in updates if u.user_id is not None}
        usernames = auth_service.get_usernames_by_ids(owner_ids)

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Weekly Project Updates"
        headers = ["Team/Project", "Week Start", "Week End", "Owner", "Status",
                   "Achievements", "Initiatives", "Next Week's Plan"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = openpyxl.styles.Font(bold=True)
        for row_idx, upd in enumerate(updates, 2):
            owner = usernames.get(upd.user_id, "") if upd.user_id else ""
            ws.cell(row=row_idx, column=1, value=upd.team_project)
            ws.cell(row=row_idx, column=2, value=str(upd.start_of_week))
            ws.cell(row=row_idx, column=3, value=str(upd.end_of_week))
            ws.cell(row=row_idx, column=4, value=owner)
            ws.cell(row=row_idx, column=5, value=upd.approval_status)
            ws.cell(row=row_idx, column=6, value=_strip_html(upd.achievements))
            ws.cell(row=row_idx, column=7, value=_strip_html(upd.initiatives))
            ws.cell(row=row_idx, column=8, value=_strip_html(upd.next_weeks_plan))
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        return Response(
            content=output.read(),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": "attachment; filename=weekly_project_updates.xlsx"},
        )

    @router.get("/export/ppt")
    def export_ppt(current_user: dict = Depends(get_current_user)) -> Response:
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Inches, Pt  # type: ignore[import]

        if current_user.get("role") != "DU_HEAD":
            raise HTTPException(status.HTTP_403_FORBIDDEN, "Only DU Head can export reports.")
        updates = list_updates_service.execute()
        owner_ids = {u.user_id for u in updates if u.user_id is not None}
        usernames = auth_service.get_usernames_by_ids(owner_ids)

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Weekly Project Updates"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Generated {datetime.utcnow().strftime('%d %B %Y')}"

        for upd in updates:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
            owner = usernames.get(upd.user_id, "") if upd.user_id else ""

            tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"{upd.team_project}  |  {upd.start_of_week} – {upd.end_of_week}"
            p.font.size = Pt(18)
            p.font.bold = True

            tb2 = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.4))
            tf2 = tb2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = f"Status: {upd.approval_status}   |   Submitted by: {owner or 'Unknown'}"
            p2.font.size = Pt(10)

            for i, (label, content) in enumerate([
                ("Achievements", upd.achievements),
                ("Initiatives", upd.initiatives),
                ("Next Week's Plan", upd.next_weeks_plan),
            ]):
                tb3 = slide.shapes.add_textbox(Inches(0.4 + i * 4.3), Inches(1.6), Inches(4.1), Inches(5.6))
                tf3 = tb3.text_frame
                tf3.word_wrap = True
                ph = tf3.paragraphs[0]
                ph.text = label
                ph.font.bold = True
                ph.font.size = Pt(11)
                pc = tf3.add_paragraph()
                pc.text = _strip_html(content)
                pc.font.size = Pt(9)

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return Response(
            content=output.read(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=weekly_project_updates.pptx"},
        )

    @router.get("/check")
    def check_report(
        start_of_week: date,
        team_project: str,
        current_user: dict = Depends(get_current_user),
    ) -> dict[str, object]:
        """Return the existing report for this week+team (any owner), or null."""
        user_id = UUID(current_user["sub"])
        existing = check_existing.execute(user_id, start_of_week, team_project)
        if existing is None:
            return {"exists": False, "report": None}
        owner_username = None
        if existing.user_id is not None:
            usernames = auth_service.get_usernames_by_ids({existing.user_id})
            owner_username = usernames.get(existing.user_id)
        return {"exists": True, "report": serialize(existing, owner_username)}

    @router.get("/{update_id}")
    def read_update(
        update_id: UUID,
        current_user: dict = Depends(get_current_user),
    ) -> dict[str, object]:
        user_id = UUID(current_user["sub"])
        user_role = current_user.get("role")
        update = get_update.execute(update_id)
        if update is None:
            raise HTTPException(status.HTTP_404_NOT_FOUND, "Project update not found.")

        if user_role in ("APP_ADMIN", "DU_HEAD", "TEAM_MANAGER"):
            pass  # full read access
        elif user_role == "TEAM_LEAD":
            delegated_manager_id = auth_service.get_active_delegation_for_delegate(user_id)
            if delegated_manager_id:
                pass  # acting as delegate → full access
            elif update.user_id != user_id:
                raise HTTPException(status.HTTP_403_FORBIDDEN, "Not authorized to view this update.")
        else:
            raise HTTPException(status.HTTP_403_FORBIDDEN, "Invalid role.")

        owner_username = None
        if update.user_id is not None:
            owner_username = auth_service.get_usernames_by_ids({update.user_id}).get(update.user_id)
        return serialize(update, owner_username)

    @router.get("")
    def list_updates(
        current_user: dict = Depends(get_current_user),
    ) -> list[dict[str, object]]:
        user_id = UUID(current_user["sub"])
        user_role = current_user.get("role")

        if user_role in ("APP_ADMIN", "DU_HEAD", "TEAM_MANAGER"):
            # All three roles see every report
            updates = list_updates_service.execute()
        elif user_role == "TEAM_LEAD":
            delegated_manager_id = auth_service.get_active_delegation_for_delegate(user_id)
            if delegated_manager_id:
                # Acting as delegate → same full visibility as a manager
                updates = list_updates_service.execute()
            else:
                updates = list_updates_service.execute({user_id})
        else:
            raise HTTPException(status.HTTP_403_FORBIDDEN, "Invalid role.")

        owner_ids = {upd.user_id for upd in updates if upd.user_id is not None}
        usernames = auth_service.get_usernames_by_ids(owner_ids)
        return [serialize(upd, usernames.get(upd.user_id)) for upd in updates]

    @router.put("/{update_id}")
    def update_project_update(
        update_id: UUID,
        request: UpdateProjectUpdateRequest,
        current_user: dict = Depends(get_current_user),
    ) -> dict[str, object]:
        user_id = UUID(current_user["sub"])
        user_role = current_user.get("role")

        if user_role == "DU_HEAD":
            raise HTTPException(status.HTTP_403_FORBIDDEN, "DU Head role cannot edit reports.")

        existing = get_update.execute(update_id)
        if existing is None:
            raise HTTPException(status.HTTP_404_NOT_FOUND, "Project update not found.")
        if existing.approval_status == "APPROVED" and user_role != "APP_ADMIN":
            raise HTTPException(status.HTTP_403_FORBIDDEN, "This report has been approved and cannot be modified.")
        if user_role == "TEAM_LEAD":
            delegated_manager_id = auth_service.get_active_delegation_for_delegate(user_id)
            if not delegated_manager_id and existing.user_id != user_id:
                raise HTTPException(status.HTTP_403_FORBIDDEN, "Only the report owner can edit this update.")
        # TEAM_MANAGER and APP_ADMIN can edit any unapproved report – no further restriction

        try:
            update = update_service.execute(
                update_id,
                user_id,
                user_role,
                UpdateCommand(
                    start_of_week=request.start_of_week,
                    end_of_week=request.end_of_week,
                    team_project=request.team_project,
                    achievements=request.achievements,
                    initiatives=request.initiatives,
                    next_weeks_plan=request.next_weeks_plan,
                ),
            )
        except DomainValidationError as error:
            raise HTTPException(status.HTTP_422_UNPROCESSABLE_ENTITY, str(error)) from error
        if update is None:
            raise HTTPException(status.HTTP_403_FORBIDDEN, "Update could not be applied.")

        owner_username = None
        if update.user_id is not None:
            owner_username = auth_service.get_usernames_by_ids({update.user_id}).get(update.user_id)
        return serialize(update, owner_username)

    @router.post("/{update_id}/approve", status_code=status.HTTP_200_OK)
    def approve_update(
        update_id: UUID,
        current_user: dict = Depends(get_current_user),
    ) -> dict[str, object]:
        user_id = UUID(current_user["sub"])
        user_role = current_user.get("role")

        if user_role not in ("TEAM_MANAGER", "APP_ADMIN", "TEAM_LEAD"):
            raise HTTPException(status.HTTP_403_FORBIDDEN, "Only Team Managers can approve reports.")

        existing = get_update.execute(update_id)
        if existing is None:
            raise HTTPException(status.HTTP_404_NOT_FOUND, "Project update not found.")

        # TEAM_MANAGER (or a delegate acting as one) can approve any unapproved report
        if user_role == "TEAM_LEAD":
            # Only allowed if this user is an active delegate
            delegated_manager_id = auth_service.get_active_delegation_for_delegate(user_id)
            if not delegated_manager_id:
                raise HTTPException(status.HTTP_403_FORBIDDEN, "Only Team Managers can approve reports.")

        update = approve_service.execute(update_id, user_id)
        if update is None:
            raise HTTPException(status.HTTP_404_NOT_FOUND, "Project update not found.")

        owner_username = None
        if update.user_id is not None:
            owner_username = auth_service.get_usernames_by_ids({update.user_id}).get(update.user_id)
        return serialize(update, owner_username)

    return router

